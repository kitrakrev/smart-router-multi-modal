#!/usr/bin/env python3
"""Generate slides.pptx from the HTML presentation content."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Colors
BG = RGBColor(0x0a, 0x0e, 0x17)
SURFACE = RGBColor(0x11, 0x18, 0x27)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT = RGBColor(0xE2, 0xE8, 0xF0)
DIM = RGBColor(0x94, 0xA3, 0xB8)
BLUE = RGBColor(0x3B, 0x82, 0xF6)
PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
GREEN = RGBColor(0x22, 0xC5, 0x5E)
RED = RGBColor(0xEF, 0x44, 0x44)
ORANGE = RGBColor(0xF9, 0x73, 0x16)
CYAN = RGBColor(0x06, 0xB6, 0xD4)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT

# Use blank layout
blank_layout = prs.slide_layouts[6]


def set_slide_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_shape(slide, text, top=Inches(0.8), font_size=Pt(36)):
    left = Inches(1)
    width = SLIDE_WIDTH - Inches(2)
    height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.bold = True
    p.font.color.rgb = BLUE
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    return txBox


def add_rule(slide, top=Inches(1.7)):
    left = Inches(3.5)
    width = Inches(6.333)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLUE
    shape.line.fill.background()
    return shape


def add_body_box(slide, top=Inches(2.1), left=Inches(1.2), width=None, height=None, font_size=Pt(16)):
    if width is None:
        width = SLIDE_WIDTH - Inches(2.4)
    if height is None:
        height = SLIDE_HEIGHT - top - Inches(0.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    return tf


def add_run(paragraph, text, color=DIM, bold=False, font_name="Consolas", font_size=Pt(15)):
    run = paragraph.add_run()
    run.text = text
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font_name
    run.font.size = font_size
    return run


def new_para(tf):
    p = tf.add_paragraph()
    return p


# ============ SLIDE 1: Title ============
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)

txBox = slide.shapes.add_textbox(Inches(1), Inches(2), SLIDE_WIDTH - Inches(2), Inches(1.5))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
add_run(p, "Smart Router for Multi-Modal LLMs", BLUE, True, "Calibri", Pt(44))

add_rule(slide, Inches(3.5))

txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3.9), SLIDE_WIDTH - Inches(2), Inches(1))
tf2 = txBox2.text_frame
tf2.word_wrap = True
p2 = tf2.paragraphs[0]
p2.alignment = PP_ALIGN.CENTER
add_run(p2, "Intelligent routing with vision, tools, and real-time observability", DIM, False, "Calibri", Pt(22))

txBox3 = slide.shapes.add_textbox(Inches(1), Inches(5.2), SLIDE_WIDTH - Inches(2), Inches(1))
tf3 = txBox3.text_frame
tf3.word_wrap = True
p3 = tf3.paragraphs[0]
p3.alignment = PP_ALIGN.CENTER
add_run(p3, "A10 Networks Hackathon \u2014 April 2026\nTeam: Karthik Raja", BLUE, False, "Calibri", Pt(18))

# ============ SLIDE 2: The Problem ============
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
add_title_shape(slide, "LLM Costs Are Exploding")
add_rule(slide)

tf = add_body_box(slide)
p = tf.paragraphs[0]
add_run(p, '"What is 2+2?"', WHITE, True)
add_run(p, '   \u2192 GPT-4o costs $0.01', RED)
add_run(p, '     (overkill)', DIM)

p = new_para(tf)
add_run(p, '"Prove Riemann"', WHITE, True)
add_run(p, '  \u2192 GPT-4o costs $0.05', GREEN)
add_run(p, '     (justified)', DIM)

p = new_para(tf); p.space_before = Pt(14)
add_run(p, '80% of queries are simple \u2192 paying 10-50x too much', ORANGE, True)

p = new_para(tf); p.space_before = Pt(20)
add_run(p, 'Without router:  ', WHITE, True)
add_run(p, '$12.00', RED, True)
add_run(p, ' per 1K queries', DIM)

p = new_para(tf)
add_run(p, 'With router:     ', WHITE, True)
add_run(p, ' $0.25', GREEN, True)
add_run(p, ' per 1K queries  \u2192  ', DIM)
add_run(p, '98% savings', GREEN, True)

# ============ SLIDE 3: What We Built ============
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
add_title_shape(slide, "Smart Multi-Modal LLM Router")
add_rule(slide)

tf = add_body_box(slide, font_size=Pt(14))
p = tf.paragraphs[0]
add_run(p, "Query \u2192 [Signals] \u2192 [Decision Matcher] \u2192 [Budget Selection] \u2192 Best Model", WHITE, True, "Consolas", Pt(16))

p = new_para(tf); p.space_before = Pt(8)
add_run(p, "  |        |              |                      |", BLUE, False, "Consolas", Pt(14))
p = new_para(tf)
add_run(p, "  |    ", BLUE, False, "Consolas", Pt(14))
add_run(p, "9 signals", CYAN, False, "Consolas", Pt(14))
add_run(p, "      ", BLUE, False, "Consolas", Pt(14))
add_run(p, "Embedding cosine", PURPLE, False, "Consolas", Pt(14))
add_run(p, "      ", BLUE, False, "Consolas", Pt(14))
add_run(p, "Cheapest model", GREEN, False, "Consolas", Pt(14))

p = new_para(tf)
add_run(p, "  |    ", BLUE, False, "Consolas", Pt(14))
add_run(p, "in parallel", CYAN, False, "Consolas", Pt(14))
add_run(p, "    ", BLUE, False, "Consolas", Pt(14))
add_run(p, "to decision exemplars", PURPLE, False, "Consolas", Pt(14))
add_run(p, "  ", BLUE, False, "Consolas", Pt(14))
add_run(p, "with required", GREEN, False, "Consolas", Pt(14))

p = new_para(tf)
add_run(p, "  |    ", BLUE, False, "Consolas", Pt(14))
add_run(p, "(10ms)", CYAN, False, "Consolas", Pt(14))
add_run(p, "         ", BLUE, False, "Consolas", Pt(14))
add_run(p, "(vLLM-SR style)", PURPLE, False, "Consolas", Pt(14))
add_run(p, "       ", BLUE, False, "Consolas", Pt(14))
add_run(p, "capabilities", GREEN, False, "Consolas", Pt(14))

p = new_para(tf); p.space_before = Pt(8)
add_run(p, "  |\u2014 ", BLUE, False, "Consolas", Pt(14))
add_run(p, "Guardrails:", ORANGE, True, "Consolas", Pt(14))
add_run(p, " jailbreak + PII detection", DIM, False, "Consolas", Pt(14))

p = new_para(tf)
add_run(p, "  |\u2014 ", BLUE, False, "Consolas", Pt(14))
add_run(p, "Vision:", CYAN, True, "Consolas", Pt(14))
add_run(p, " image-aware routing", DIM, False, "Consolas", Pt(14))

p = new_para(tf)
add_run(p, "  |\u2014 ", BLUE, False, "Consolas", Pt(14))
add_run(p, "Tools:", PURPLE, True, "Consolas", Pt(14))
add_run(p, " function calling support", DIM, False, "Consolas", Pt(14))

p = new_para(tf)
add_run(p, "  |\u2014 ", BLUE, False, "Consolas", Pt(14))
add_run(p, "Adaptive:", GREEN, True, "Consolas", Pt(14))
add_run(p, " auto-disable degraded models", DIM, False, "Consolas", Pt(14))

# ============ SLIDE 4: Architecture ============
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
add_title_shape(slide, "Architecture")
add_rule(slide)

# Architecture box with dark surface bg
left = Inches(1.5)
top = Inches(2.2)
width = SLIDE_WIDTH - Inches(3)
height = Inches(4.8)
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
shape.fill.solid()
shape.fill.fore_color.rgb = SURFACE
shape.line.color.rgb = RGBColor(0x1E, 0x29, 0x3B)

tf = shape.text_frame
tf.word_wrap = False
sz = Pt(12)

lines = [
    [("+-----------------------------------------------------+", BLUE)],
    [("|", BLUE), ("                   ", None), ("FastAPI Server", WHITE), ("                     ", None), ("|", BLUE)],
    [("|", BLUE), ("              ", None), ("OpenAI-Compatible API", CYAN), ("                   ", None), ("|", BLUE)],
    [("+-----------------------------------------------------+", BLUE)],
    [("|", BLUE), ("                                                      ", None), ("|", BLUE)],
    [("|", BLUE), ("  Request \u2192 ", DIM), ("Guardrails", RED), (" \u2192 ", DIM), ("Signals", CYAN), (" \u2192 ", DIM), ("Decision", PURPLE), (" \u2192 ", DIM), ("Model", GREEN), ("  ", None), ("|", BLUE)],
    [("|", BLUE), ("              ", None), ("|", RED), ("            ", None), ("|", CYAN), ("          ", None), ("|", PURPLE), ("          ", None), ("|", GREEN), ("     ", None), ("|", BLUE)],
    [("|", BLUE), ("          ", None), ("jailbreak", RED), ("    ", None), ("keyword", CYAN), ("     ", None), ("embedding", PURPLE), ("   ", None), ("budget", GREEN), ("  ", None), ("|", BLUE)],
    [("|", BLUE), ("          ", None), ("PII", RED), ("         ", None), ("domain", CYAN), ("      ", None), ("match to", PURPLE), ("    ", None), ("aware", GREEN), ("   ", None), ("|", BLUE)],
    [("|", BLUE), ("                      ", None), ("complexity", CYAN), ("  ", None), ("exemplars", PURPLE), ("   ", None), ("select", GREEN), ("  ", None), ("|", BLUE)],
    [("|", BLUE), ("                      ", None), ("language", CYAN), ("                         ", None), ("|", BLUE)],
    [("|", BLUE), ("                      ", None), ("vision", CYAN), ("                           ", None), ("|", BLUE)],
    [("|", BLUE), ("                      ", None), ("tools", CYAN), ("                            ", None), ("|", BLUE)],
    [("|", BLUE), ("                      ", None), ("modality", CYAN), ("                         ", None), ("|", BLUE)],
    [("|", BLUE), ("                                                      ", None), ("|", BLUE)],
    [("+-----------------------------------------------------+", BLUE)],
    [("|", BLUE), (" ", None), ("Model Registry", WHITE), ("    ", None), ("|", BLUE), ("  ", None), ("Adaptive Routing", WHITE), ("  ", None), ("|", BLUE), (" ", None), ("Dashboard", WHITE), ("  ", None), ("|", BLUE)],
    [("|", BLUE), (" Add/remove/discover ", DIM), ("|", BLUE), ("  Auto-disable/recover ", DIM), ("|", BLUE), (" 5-tab UI  ", DIM), ("|", BLUE)],
    [("|", BLUE), (" at runtime          ", DIM), ("|", BLUE), ("  EMA stats           ", DIM), ("|", BLUE), (" WebSocket ", DIM), ("|", BLUE)],
    [("+-----------------------------------------------------+", BLUE)],
]

for i, line_parts in enumerate(lines):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = new_para(tf)
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    for text, color in line_parts:
        if color is None:
            add_run(p, text, DIM, False, "Consolas", sz)
        else:
            add_run(p, text, color, False, "Consolas", sz)

# ============ SLIDE 5: Why Embedding Match ============
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
add_title_shape(slide, "Why Embedding Match, Not Rules?")
add_rule(slide)

# Left column
tf_left = add_body_box(slide, Inches(2.3), Inches(1), Inches(5.2), Inches(4.5))
p = tf_left.paragraphs[0]
add_run(p, "Rule-Based Routing (Brittle)", RED, True, "Calibri", Pt(18))

p = new_para(tf_left); p.space_before = Pt(14)
add_run(p, 'if ', DIM, False, "Consolas", Pt(14))
add_run(p, '"prove"', ORANGE)
add_run(p, ' in query \u2192 math', DIM)

p = new_para(tf_left)
add_run(p, 'if ', DIM, False, "Consolas", Pt(14))
add_run(p, '"code"', ORANGE)
add_run(p, ' in query \u2192 code', DIM)

p = new_para(tf_left)
add_run(p, 'if ', DIM, False, "Consolas", Pt(14))
add_run(p, '"write"', ORANGE)
add_run(p, ' in query \u2192 ???', DIM)

p = new_para(tf_left); p.space_before = Pt(12)
add_run(p, 'Breaks on: ', RED, True)
add_run(p, '"analyze this\nfunction\'s time complexity"', DIM)

p = new_para(tf_left); p.space_before = Pt(8)
add_run(p, 'Manual rules for every edge case', DIM)

# Right column
tf_right = add_body_box(slide, Inches(2.3), Inches(7), Inches(5.2), Inches(4.5))
p = tf_right.paragraphs[0]
add_run(p, "Embedding Match (Robust)", GREEN, True, "Calibri", Pt(18))

p = new_para(tf_right); p.space_before = Pt(14)
add_run(p, 'query embed \u2192 ', DIM)
add_run(p, 'cosine sim', PURPLE)

p = new_para(tf_right)
add_run(p, 'to decision exemplar centroids', DIM)

p = new_para(tf_right); p.space_before = Pt(12)
add_run(p, '"analyze complexity" \u2192 ', DIM)
add_run(p, 'code (0.48)', GREEN)

p = new_para(tf_right)
add_run(p, 'NOT math, NOT creative', DIM)

p = new_para(tf_right)
add_run(p, 'Handles ambiguity naturally', GREEN, True)

p = new_para(tf_right); p.space_before = Pt(12)
add_run(p, 'Just add exemplars to YAML', DIM)

# ============ SLIDE 6: Budget-Aware Selection ============
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
add_title_shape(slide, "Budget-Aware Selection")
add_rule(slide)

tf = add_body_box(slide)
p = tf.paragraphs[0]
add_run(p, 'Not Just "Which Model" \u2014 "Cheapest Correct Model"', WHITE, True, "Calibri", Pt(20))

p = new_para(tf); p.space_before = Pt(18)
add_run(p, 'Decision: ', DIM)
add_run(p, '"This is a code task"', PURPLE)

p = new_para(tf)
add_run(p, '  require: ', DIM)
add_run(p, '[code]', CYAN)

p = new_para(tf)
add_run(p, '  strategy: ', DIM)
add_run(p, 'cheapest_capable', BLUE)

p = new_para(tf); p.space_before = Pt(16)
add_run(p, 'Available models:', DIM)

p = new_para(tf); p.space_before = Pt(6)
add_run(p, '  \u2022 GPT-4o       ', ORANGE)
add_run(p, '[code] \u2713', GREEN)
add_run(p, '  $10.00/1K', RED)
add_run(p, '   \u2190 expensive', DIM)

p = new_para(tf)
add_run(p, '  \u2022 Claude Sonnet ', ORANGE)
add_run(p, '[code] \u2713', GREEN)
add_run(p, '  $3.00/1K', ORANGE)

p = new_para(tf)
add_run(p, '  \u2022 Gemini Flash  ', ORANGE)
add_run(p, '[code] \u2713', GREEN)
add_run(p, '  $0.40/1K', GREEN)
add_run(p, '   \u2190 ', DIM)
add_run(p, 'PICKED \u2713', GREEN, True)

p = new_para(tf)
add_run(p, '  \u2022 Qwen-3B      ', ORANGE)
add_run(p, '[text] \u2717', DIM)
add_run(p, '              \u2190 no code capability', DIM)

p = new_para(tf); p.space_before = Pt(16)
add_run(p, '3 strategies: ', DIM)
add_run(p, 'cheapest_capable', BLUE)
add_run(p, ' | ', DIM)
add_run(p, 'quality_first', PURPLE)
add_run(p, ' | ', DIM)
add_run(p, 'balanced', CYAN)

# ============ SLIDE 7: What's New ============
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
add_title_shape(slide, "What's New (vs vLLM-SR)")
add_rule(slide)

# Left col
tf_left = add_body_box(slide, Inches(2.3), Inches(1), Inches(5.2), Inches(4.8))
p = tf_left.paragraphs[0]
add_run(p, "vLLM Semantic Router", BLUE, True, "Calibri", Pt(18))

items_left = [
    ("\u2713 14 signals", GREEN),
    ("\u2713 Embedding routing", GREEN),
    ("\u2713 Jailbreak detection", GREEN),
    ("\u2713 Semantic cache", GREEN),
    ("\u2713 14 selection algorithms", GREEN),
    ("", None),
    ("\u2717 Vision content analysis", RED),
    ("\u2717 Tool execution", RED),
    ("\u2717 Runtime decision path UI", RED),
    ("\u2717 Adaptive model health", RED),
    ("\u2717 Budget optimization", RED),
]
for text, color in items_left:
    p = new_para(tf_left)
    if text:
        p.space_before = Pt(4)
        add_run(p, text, color)
    else:
        p.space_before = Pt(8)

# Right col
tf_right = add_body_box(slide, Inches(2.3), Inches(7), Inches(5.2), Inches(4.8))
p = tf_right.paragraphs[0]
add_run(p, "Our Smart Router", GREEN, True, "Calibri", Pt(18))

p = new_para(tf_right); p.space_before = Pt(8)
add_run(p, "\u2713 9 signals + ", GREEN)
add_run(p, "3 NEW:", CYAN, True)

for item in ["Vision (CLIP image routing)", "Tool awareness", "Modality detection"]:
    p = new_para(tf_right); p.space_before = Pt(2)
    add_run(p, "   " + item, CYAN)

new_items = [
    "Budget-aware selection",
    "Adaptive routing (auto-disable)",
    "Dynamic model registry",
    "Session/user tracking",
    "Real-time trace dashboard",
    "Hot-reload config API",
    "Simulate degradation (demo)",
]
for item in new_items:
    p = new_para(tf_right); p.space_before = Pt(4)
    add_run(p, "\u2713 NEW: ", GREEN, True)
    add_run(p, item, GREEN)

# ============ SLIDE 8: Benchmark Results ============
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
add_title_shape(slide, "Benchmark Results")
add_rule(slide)

tf = add_body_box(slide, Inches(2.1))
p = tf.paragraphs[0]
add_run(p, "RouterArena (8,400 text queries)", WHITE, True, "Calibri", Pt(18))

p = new_para(tf); p.space_before = Pt(12)
add_run(p, "                  Accuracy    Latency    Cost/1K", DIM, False, "Consolas", Pt(14))
p = new_para(tf); p.space_before = Pt(4)
add_run(p, "  Our Router      ", GREEN, True, "Consolas", Pt(14))
add_run(p, "98.1%", GREEN, True, "Consolas", Pt(14))
add_run(p, "       ", DIM, False, "Consolas", Pt(14))
add_run(p, "10.6ms", GREEN, True, "Consolas", Pt(14))
add_run(p, "     ", DIM, False, "Consolas", Pt(14))
add_run(p, "$233", ORANGE, False, "Consolas", Pt(14))

p = new_para(tf); p.space_before = Pt(2)
add_run(p, "  vLLM-SR         ", DIM, False, "Consolas", Pt(14))
add_run(p, "66.5%", RED, False, "Consolas", Pt(14))
add_run(p, "       ", DIM, False, "Consolas", Pt(14))
add_run(p, "243ms", RED, False, "Consolas", Pt(14))
add_run(p, "      ", DIM, False, "Consolas", Pt(14))
add_run(p, "$0.06", GREEN, False, "Consolas", Pt(14))

p = new_para(tf); p.space_before = Pt(2)
add_run(p, "  Always GPT-4o   ", DIM, False, "Consolas", Pt(14))
add_run(p, "\u2014", DIM, False, "Consolas", Pt(14))
add_run(p, "             ", DIM, False, "Consolas", Pt(14))
add_run(p, "\u2014", DIM, False, "Consolas", Pt(14))
add_run(p, "          ", DIM, False, "Consolas", Pt(14))
add_run(p, "$12.00", RED, False, "Consolas", Pt(14))

p = new_para(tf); p.space_before = Pt(24)
add_run(p, "VL-RouterBench (30K vision queries)", WHITE, True, "Calibri", Pt(18))

p = new_para(tf); p.space_before = Pt(12)
add_run(p, "                  Accuracy    Rank Score", DIM, False, "Consolas", Pt(14))
p = new_para(tf); p.space_before = Pt(4)
add_run(p, "  Our Router      ", GREEN, True, "Consolas", Pt(14))
add_run(p, "83.8%", GREEN, True, "Consolas", Pt(14))
add_run(p, "       ", DIM, False, "Consolas", Pt(14))
add_run(p, "0.80", GREEN, True, "Consolas", Pt(14))

p = new_para(tf); p.space_before = Pt(2)
add_run(p, "  Always GPT-4o   84.4%       ", DIM, False, "Consolas", Pt(14))
add_run(p, "0.48", RED, False, "Consolas", Pt(14))
add_run(p, " (expensive)", DIM, False, "Consolas", Pt(14))

p = new_para(tf); p.space_before = Pt(2)
add_run(p, "  Random          ", DIM, False, "Consolas", Pt(14))
add_run(p, "58.6%", RED, False, "Consolas", Pt(14))
add_run(p, "       0.55", DIM, False, "Consolas", Pt(14))

p = new_para(tf); p.space_before = Pt(16)
add_run(p, "Custom embeddings trained on LMSYS 55K (pushed to HuggingFace)", DIM, False, "Consolas", Pt(12))

# ============ SLIDE 9: Live Demo ============
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
add_title_shape(slide, "Live Demo")
add_rule(slide)

tf = add_body_box(slide)
p = tf.paragraphs[0]
add_run(p, "http://localhost:8000/dashboard", BLUE, False, "Consolas", Pt(18))

p = new_para(tf); p.space_before = Pt(20)
add_run(p, "Demo Flow:", WHITE, True, "Calibri", Pt(20))

demos = [
    ("1.", " Send simple query     \u2192 routes to cheap model ", "qwen-3b", GREEN),
    ("2.", " Send math proof       \u2192 routes to reasoning model ", "claude", PURPLE),
    ("3.", " Send code query       \u2192 routes to code model ", "gemini-flash", BLUE),
    ("4.", " Send jailbreak        \u2192 ", "BLOCKED by guardrails", RED),
    ("5.", " Simulate degradation  \u2192 watch ", "auto-disable", ORANGE),
    ("6.", " Add new model via API \u2192 immediately available", None, None),
    ("7.", " Show session trace    \u2192 full pipeline visualization", None, None),
]
for num, desc, highlight, hcolor in demos:
    p = new_para(tf); p.space_before = Pt(6)
    add_run(p, num, CYAN, True)
    if highlight:
        add_run(p, desc, DIM)
        add_run(p, highlight, hcolor, True)
    else:
        add_run(p, desc, DIM)

# ============ SLIDE 10: Technical Decisions ============
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
add_title_shape(slide, "Technical Decisions")
add_rule(slide)

tf = add_body_box(slide)

sections = [
    ("Python over Go/Rust", [
        ("\u2192 ", DIM, False), ("3,300 lines", GREEN, True), (" vs ", DIM, False), ("654K lines", RED, True), (" (vLLM-SR)", DIM, False),
    ]),
    (None, [("\u2192 Same accuracy, faster to iterate", DIM, False)]),
    (None, [("\u2192 10ms latency (good enough for routing)", DIM, False)]),
    ("MiniLM-L6 over mmBERT-32K", [
        ("\u2192 ", DIM, False), ("22M", GREEN, True), (" vs ", DIM, False), ("350M", RED, True), (" params", DIM, False),
    ]),
    (None, [("\u2192 ", DIM, False), ("10ms", GREEN, True), (" vs ", DIM, False), ("243ms", RED, True), (" per query", DIM, False)]),
    (None, [("\u2192 ", DIM, False), ("98%", GREEN, True), (" vs ", DIM, False), ("66%", RED, True), (" accuracy (better with good exemplars)", DIM, False)]),
    ("No model names in training", [
        ("\u2192 Trained on ", DIM, False), ("task profiles", PURPLE, True), (", not \"pick GPT-4\"", DIM, False),
    ]),
    (None, [("\u2192 Add new model tomorrow \u2192 works without retraining", DIM, False)]),
    (None, [("\u2192 Budget layer maps profiles \u2192 cheapest capable model", DIM, False)]),
    ("Embedding match over classifier", [
        ("\u2192 Zero training needed \u2014 add exemplars in YAML", DIM, False),
    ]),
    (None, [("\u2192 ", DIM, False), ("98.1% accuracy", GREEN, True), (" without any training data", DIM, False)]),
]

first = True
for title, parts in sections:
    p = new_para(tf) if not first else tf.paragraphs[0]
    first = False
    if title:
        p.space_before = Pt(14)
        add_run(p, title, WHITE, True, "Calibri", Pt(16))
        p = new_para(tf)
    p.space_before = Pt(2)
    for text, color, bold in parts:
        add_run(p, text, color, bold, "Consolas", Pt(14))

# ============ SLIDE 11: Honest Limitations ============
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
add_title_shape(slide, "What We Got Wrong & Learned")
add_rule(slide)

tf = add_body_box(slide)

limitations = [
    ("RouterArena 98.1% is inflated", [
        "gpt-4o-mini has ground truth for all queries",
        "Routing away from it = \"wrong\" by default",
    ], "Real accuracy on 3-model overlap: ", "80.9%", ORANGE),
    ("Removing thresholds dropped accuracy to 51%", [
        "\"What is 2+2\" \u2192 complex_reasoning (oops)",
        "Need soft thresholds, not hard rules",
    ], None, None, None),
    ("Fine-tuning on benchmark data = data leakage", [
        "Caught it, switched to LMSYS for training",
    ], None, None, None),
    ("Keyword fallback was silently active", [
        "Embeddings weren't loading, nobody noticed",
    ], "Added ", "loud failure mode", GREEN),
]

first = True
for title, bullets, extra_text, extra_highlight, extra_color in limitations:
    p = new_para(tf) if not first else tf.paragraphs[0]
    first = False
    p.space_before = Pt(14)
    add_run(p, "\u2717 ", RED, True)
    add_run(p, title, WHITE, True, "Calibri", Pt(16))

    for b in bullets:
        p = new_para(tf); p.space_before = Pt(2)
        add_run(p, "  \u2192 " + b, DIM)

    if extra_text:
        p = new_para(tf); p.space_before = Pt(2)
        add_run(p, "  \u2192 " + extra_text, DIM)
        add_run(p, extra_highlight, extra_color, True)
        if extra_text.startswith("Added"):
            add_run(p, " \u2014 no silent fallbacks", DIM)

# ============ SLIDE 12: A10 Connection ============
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
add_title_shape(slide, "How This Fits A10's Mission")
add_rule(slide)

tf = add_body_box(slide)
p = tf.paragraphs[0]
add_run(p, "A10 AI Firewall (guardrails)", WHITE, True, "Calibri", Pt(18))
add_run(p, "     ", DIM)
add_run(p, "Our Router (optimization)", WHITE, True, "Calibri", Pt(18))

p = new_para(tf); p.space_before = Pt(16)
add_run(p, '"Is this prompt SAFE?"', RED, True)
add_run(p, '           ', DIM)
add_run(p, '"Which model HANDLES this best?"', GREEN, True)

p = new_para(tf); p.space_before = Pt(24)
add_run(p, 'Together:', WHITE, True, "Calibri", Pt(20))

p = new_para(tf); p.space_before = Pt(12)
add_run(p, '  Prompt \u2192 ', DIM)
add_run(p, 'A10 Guardrail', RED, True)
add_run(p, ' \u2192 ', DIM)
add_run(p, 'Smart Router', BLUE, True)
add_run(p, ' \u2192 ', DIM)
add_run(p, 'Best LLM', GREEN, True)
add_run(p, ' \u2192 Response', DIM)

p = new_para(tf); p.space_before = Pt(4)
add_run(p, '           ', DIM)
add_run(p, '(block threats)', RED)
add_run(p, '  ', DIM)
add_run(p, '(optimize cost)', BLUE)
add_run(p, '  ', DIM)
add_run(p, '(generate)', GREEN)

p = new_para(tf); p.space_before = Pt(24)
add_run(p, 'From ', DIM)
add_run(p, 'ATEN-tion', PURPLE, True)
add_run(p, ' (detect) to ', DIM)
add_run(p, 'In-TEN-tion', PURPLE, True)
add_run(p, ' (act intelligently)', DIM)

# ============ SLIDE 13: Repo & Links ============
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
add_title_shape(slide, "Repo & Links")
add_rule(slide)

tf = add_body_box(slide)
p = tf.paragraphs[0]
add_run(p, "github.com/kitrakrev/smart-router-multi-modal", BLUE, False, "Consolas", Pt(18))

tree = [
    ("\u251c\u2500\u2500", "server.py", "          FastAPI, OpenAI-compatible"),
    ("\u251c\u2500\u2500", "signals.py", "         9 signal extractors"),
    ("\u251c\u2500\u2500", "router.py", "          Decision embedding matcher"),
    ("\u251c\u2500\u2500", "models.py", "          Dynamic model registry"),
    ("\u251c\u2500\u2500", "tools.py", "           11 security tools"),
    ("\u251c\u2500\u2500", "index.html", "         5-tab live dashboard"),
    ("\u251c\u2500\u2500", "benchmarks/", "        RouterArena + VL-RouterBench"),
    ("\u2514\u2500\u2500", "train_routing_embeddings.py", ""),
]
for prefix, name, desc in tree:
    p = new_para(tf); p.space_before = Pt(4)
    add_run(p, prefix + " ", DIM, False, "Consolas", Pt(14))
    add_run(p, name, WHITE, True, "Consolas", Pt(14))
    if desc:
        add_run(p, desc, DIM, False, "Consolas", Pt(14))

p = new_para(tf); p.space_before = Pt(20)
add_run(p, "HuggingFace: ", PURPLE, True, "Consolas", Pt(15))
add_run(p, "kitrakrev/smart-router-embeddings", DIM, False, "Consolas", Pt(15))

p = new_para(tf); p.space_before = Pt(4)
add_run(p, "Dashboard:   ", BLUE, True, "Consolas", Pt(15))
add_run(p, "http://localhost:8000/dashboard", DIM, False, "Consolas", Pt(15))

p = new_para(tf); p.space_before = Pt(24)
add_run(p, "Thank you!", DIM, False, "Calibri", Pt(22))

# Save
output_path = os.path.join(os.path.dirname(__file__), "slides.pptx")
prs.save(output_path)
print(f"Saved to {output_path}")
